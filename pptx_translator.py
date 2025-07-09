"""
pptx翻訳スクリプト
日本語のスライドショーを、レイアウト等は変えずに文字だけ翻訳する
"""

import sys
import argparse
from pathlib import Path
from typing import List, Dict, Optional
import time
from pptx import Presentation
from pptx.shapes.base import BaseShape
import google.generativeai as genai


class PPTXTranslator:
    def __init__(self, source_lang: str = "ja", target_lang: str = "en", gemini_api_key: str = None):
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.translation_cache = {}
        self.gemini_api_key = gemini_api_key
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)
            self.model = genai.GenerativeModel('gemini-2.5-flash')
        
    def extract_text_from_shape(self, shape: BaseShape) -> List[str]:
        """シェイプからテキストを抽出"""
        texts = []
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text.strip())
        return texts
    
    def extract_all_texts_with_positions(self, pptx_path: str) -> Dict:
        """PPTXファイルから全テキストと位置情報を抽出"""
        prs = Presentation(pptx_path)
        text_data = {
            "texts": [],
            "positions": []
        }
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        text = paragraph.text.strip()
                        if text:
                            # run情報も保存（安全に処理）
                            runs_info = []
                            try:
                                for run in paragraph.runs:
                                    if run.text:
                                        runs_info.append({
                                            "text": run.text,
                                            "font": run.font
                                        })
                            except Exception as e:
                                # run情報の取得に失敗した場合は空のリストを使用
                                runs_info = []
                            
                            text_data["texts"].append(text)
                            text_data["positions"].append({
                                "slide_idx": slide_idx,
                                "shape_idx": shape_idx,
                                "para_idx": para_idx,
                                "original_text": text,
                                "runs_info": runs_info
                            })
        
        return text_data
    
    def translate_batch_with_gemini(self, texts: List[str]) -> List[str]:
        """Gemini APIを使用した一括翻訳"""
        if not texts:
            return []
        
        print(f"翻訳対象テキスト数: {len(texts)}")
        
        # すべてのテキストを一度に翻訳（キャッシュは使用しない）
        # これにより順序の整合性を保つ
        texts_to_translate = texts
        
        try:
            # 言語名の変換
            lang_map = {
                "ja": "Japanese",
                "en": "English",
                "ko": "Korean",
                "zh": "Chinese",
                "es": "Spanish",
                "fr": "French",
                "de": "German"
            }
            
            source_lang_name = lang_map.get(self.source_lang, self.source_lang)
            target_lang_name = lang_map.get(self.target_lang, self.target_lang)
            
            # 一括翻訳プロンプト
            texts_formatted = "\n".join([f"{i+1}. {text}" for i, text in enumerate(texts_to_translate)])
            
            prompt = f"""
Translate the following {source_lang_name} texts to {target_lang_name}.

IMPORTANT INSTRUCTIONS:
- Keep the same length and tone as the original text for each item
- Maintain the same level of formality
- If the original is concise, keep the translation concise
- If the original is detailed, keep the translation detailed
- Preserve any technical terms appropriately
- Do not add explanations or additional context
- Return only the translations, numbered exactly as shown below

Original texts:
{texts_formatted}

Translations:"""
            
            response = self.model.generate_content(prompt)
            translated_text = response.text.strip()
            
            # レスポンスから翻訳を抽出
            translations = self._parse_batch_translations(translated_text, len(texts_to_translate), texts_to_translate)
            
            # キャッシュに保存
            for original, translation in zip(texts_to_translate, translations):
                self.translation_cache[original] = translation
            
            time.sleep(0.1)  # API制限対策
            return translations
            
        except Exception as e:
            print(f"Gemini一括翻訳エラー: {e}")
            return texts
    
    def _parse_batch_translations(self, response_text: str, expected_count: int, original_texts: List[str] = None) -> List[str]:
        """バッチ翻訳のレスポンスをパース"""
        lines = response_text.strip().split('\n')
        translations = []
        
        # 番号順に翻訳を取得
        for i in range(1, expected_count + 1):
            found_translation = None
            
            for line in lines:
                line = line.strip()
                if line.startswith(f"{i}."):
                    # 番号を削除して翻訳テキストを抽出
                    try:
                        translation = line.split('.', 1)[1].strip()
                        found_translation = translation
                        break
                    except IndexError:
                        continue
            
            # 翻訳が見つからない場合は元のテキストを使用
            if found_translation is None:
                print(f"警告: 翻訳 {i} が見つかりません。元のテキストを使用します。")
                # 元のテキストを使用
                if original_texts and i <= len(original_texts):
                    found_translation = original_texts[i-1]  # 1-based indexなので-1
                else:
                    found_translation = f"[翻訳失敗-{i}]"
            
            translations.append(found_translation)
        
        # デバッグ情報を出力
        if len(translations) != expected_count:
            print(f"警告: 翻訳結果の数が期待値と異なります ({len(translations)} vs {expected_count})")
            print(f"レスポンステキスト（最初の500文字）: {response_text[:500]}...")
        
        print(f"翻訳完了: {len(translations)} 件")
        return translations
    
    def replace_text_with_format_preservation(self, paragraph, translated_text: str, runs_info: List[Dict]):
        """書式を保持してテキストを置換"""
        try:
            if not runs_info:
                # run情報がない場合は従来の方法
                paragraph.text = translated_text
                return
            
            # 元のrun数が1つの場合は単純置換
            if len(runs_info) == 1:
                paragraph.text = translated_text
                if paragraph.runs and len(paragraph.runs) > 0:
                    # 元のフォント設定を復元
                    original_font = runs_info[0].get("font")
                    if original_font:
                        current_font = paragraph.runs[0].font
                        self._copy_font_properties(original_font, current_font)
                return
            
            # 複数runの場合は比例配分で文字を分割
            total_length = sum(len(run_info.get("text", "")) for run_info in runs_info)
            if total_length == 0:
                paragraph.text = translated_text
                return
            
            # 既存のrunをクリア
            paragraph.text = ""
            
            # 翻訳テキストを元のrunの長さに比例して分割
            translated_parts = []
            start_pos = 0
            
            for i, run_info in enumerate(runs_info):
                original_length = len(run_info.get("text", ""))
                if i == len(runs_info) - 1:
                    # 最後のrunは残り全て
                    part = translated_text[start_pos:]
                else:
                    # 比例配分で長さを計算
                    part_length = int(len(translated_text) * original_length / total_length)
                    part = translated_text[start_pos:start_pos + part_length]
                    start_pos += part_length
                
                translated_parts.append(part)
            
            # 新しいrunを作成して書式を適用
            for i, (part, run_info) in enumerate(zip(translated_parts, runs_info)):
                if part:  # 空でない場合のみ追加
                    try:
                        if i == 0:
                            # 最初のrunは既存のものを使用
                            if len(paragraph.runs) > 0:
                                run = paragraph.runs[0]
                            else:
                                run = paragraph.add_run()
                        else:
                            # 新しいrunを追加
                            run = paragraph.add_run()
                        
                        run.text = part
                        
                        # フォント情報をコピー
                        original_font = run_info.get("font")
                        if original_font:
                            self._copy_font_properties(original_font, run.font)
                    except Exception as e:
                        print(f"run作成エラー (インデックス {i}): {e}")
                        continue
        
        except Exception as e:
            # 全体的なエラーの場合は単純な置換にフォールバック
            print(f"書式保持エラー、単純置換にフォールバック: {e}")
            paragraph.text = translated_text
    
    def _copy_font_properties(self, source_font, target_font):
        """フォントプロパティをコピー"""
        try:
            # フォント名をコピー
            if hasattr(source_font, 'name') and source_font.name:
                target_font.name = source_font.name
        except Exception as e:
            pass
        
        try:
            # フォントサイズをコピー
            if hasattr(source_font, 'size') and source_font.size:
                target_font.size = source_font.size
        except Exception as e:
            pass
        
        try:
            # 太字をコピー
            if hasattr(source_font, 'bold') and source_font.bold is not None:
                target_font.bold = source_font.bold
        except Exception as e:
            pass
        
        try:
            # 斜体をコピー
            if hasattr(source_font, 'italic') and source_font.italic is not None:
                target_font.italic = source_font.italic
        except Exception as e:
            pass
        
        try:
            # 下線をコピー
            if hasattr(source_font, 'underline') and source_font.underline is not None:
                target_font.underline = source_font.underline
        except Exception as e:
            pass
        
        try:
            # 色情報をコピー（安全に処理）
            if hasattr(source_font, 'color') and source_font.color:
                source_color = source_font.color
                target_color = target_font.color
                
                # RGB色の場合
                if hasattr(source_color, 'rgb') and source_color.rgb is not None:
                    target_color.rgb = source_color.rgb
                # テーマ色の場合
                elif hasattr(source_color, 'theme_color') and source_color.theme_color is not None:
                    target_color.theme_color = source_color.theme_color
                # その他の色タイプは無視
                
        except Exception as e:
            # 色情報のコピーエラーは頻繁に発生するので、無視する
            pass
    
    def translate_pptx(self, input_path: str, output_path: str) -> bool:
        """PPTXファイルを翻訳（一括処理）"""
        try:
            if not self.gemini_api_key:
                print("エラー: Gemini APIキーが設定されていません")
                return False
                
            # 全テキストと位置情報を抽出
            text_data = self.extract_all_texts_with_positions(input_path)
            texts = text_data["texts"]
            positions = text_data["positions"]
            
            if not texts:
                print("翻訳対象のテキストが見つかりません")
                return False
            
            print(f"翻訳開始: {input_path}")
            print(f"翻訳対象テキスト数: {len(texts)}")
            
            # 一括翻訳実行
            translations = self.translate_batch_with_gemini(texts)
            
            # PPTXファイルを開いて翻訳を適用
            prs = Presentation(input_path)
            
            # 安全な数でループ（小さい方を使用）
            safe_count = min(len(translations), len(positions))
            
            for i in range(safe_count):
                try:
                    translation = translations[i]
                    position = positions[i]
                    
                    slide = prs.slides[position["slide_idx"]]
                    shape = slide.shapes[position["shape_idx"]]
                    paragraph = shape.text_frame.paragraphs[position["para_idx"]]
                    
                    original_text = position["original_text"]
                    runs_info = position.get("runs_info", [])
                    
                    # デバッグ情報を出力
                    print(f"[{i+1}/{safe_count}] スライド{position['slide_idx']+1}, シェイプ{position['shape_idx']}, パラグラフ{position['para_idx']}")
                    print(f"  元テキスト: {original_text[:50]}...")
                    print(f"  翻訳結果: {translation[:50]}...")
                    
                    # 書式を保持してテキストを置換
                    self.replace_text_with_format_preservation(paragraph, translation, runs_info)
                    
                except Exception as e:
                    print(f"テキスト置換エラー (インデックス {i}): {e}")
                    continue
            
            # 処理されなかった項目があれば警告
            if safe_count < len(positions):
                print(f"警告: {len(positions) - safe_count} 個のテキストが処理されませんでした")
            
            prs.save(output_path)
            print(f"翻訳完了: {output_path}")
            return True
            
        except Exception as e:
            print(f"翻訳処理エラー: {e}")
            return False


def main():
    parser = argparse.ArgumentParser(description="PPTX翻訳スクリプト (Gemini API使用)")
    parser.add_argument("input_file", help="入力PPTXファイル")
    parser.add_argument("-o", "--output", help="出力PPTXファイル")
    parser.add_argument("-s", "--source", default="ja", help="翻訳元言語 (デフォルト: ja)")
    parser.add_argument("-t", "--target", default="en", help="翻訳先言語 (デフォルト: en)")
    parser.add_argument("--gemini-key", required=True, help="Gemini API キー")
    
    args = parser.parse_args()
    
    # 入力ファイルの確認
    if not Path(args.input_file).exists():
        print(f"エラー: ファイルが見つかりません: {args.input_file}")
        sys.exit(1)
    
    # 出力ファイル名の生成
    if args.output:
        output_file = args.output
    else:
        input_path = Path(args.input_file)
        output_file = str(input_path.parent / f"{input_path.stem}_translated.pptx")
    
    # 翻訳処理
    translator = PPTXTranslator(args.source, args.target, args.gemini_key)
    
    success = translator.translate_pptx(args.input_file, output_file)
    
    if success:
        print("翻訳が完了しました！")
    else:
        print("翻訳に失敗しました。")
        sys.exit(1)


if __name__ == "__main__":
    main()