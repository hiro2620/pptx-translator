"""
pptx翻訳スクリプト
日本語のスライドショーを、レイアウト等は変えずに文字だけ翻訳する
"""

import logging
import sys
import argparse
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Any
import shutil
import json
from dataclasses import dataclass
from pptx import Presentation
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.shapes.shapetree import SlideShapes
from google import genai


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class RunStyleInfo:
    """runのスタイル情報を格納するデータクラス"""
    text: str
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color_rgb: Optional[tuple] = None
    color_theme: Optional[int] = None
    color_brightness: Optional[float] = None

@dataclass
class TextPosition:
    """テキストの位置情報を格納するデータクラス"""
    slide_idx: int
    shape_idx: int
    para_idx: int
    original_text: str
    run_styles: List[RunStyleInfo]

@dataclass
class TableCellPosition(TextPosition):
    """テーブルセルのテキスト位置情報を格納するデータクラス"""
    table_row: int = -1
    table_col: int = -1

@dataclass
class GroupedShapePosition(TextPosition):
    """グループ化された図形のテキスト位置情報を格納するデータクラス"""
    group_path: List[int] = None  # グループ内のパス（ネストしたインデックス）
    
    def __post_init__(self) -> None:
        if self.group_path is None:
            self.group_path = []

@dataclass
class TextPositionPair:
    """テキストと位置情報のペアを格納するデータクラス"""
    text: str
    position: TextPosition

@dataclass
class PresentationData:
    """プレゼンテーション全体のデータを格納するデータクラス"""
    all_pairs: List[TextPositionPair]


class PPTXTranslator:
    def __init__(self, source_lang: str = "ja", target_lang: str = "en", model_name: str = "gemini-2.5-flash", logger: logging.Logger = logger, show_correspondence: bool = False) -> None:
        self.logger = logger
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.show_correspondence = show_correspondence
        
        # JSON Schemaを定義（翻訳後run分割と対応付きで）
        self.translation_schema = {
            "type": "object",
            "properties": {
                "translations": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "id": {"type": "integer"},
                            "translated": {"type": "string"},
                            "translated_runs": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "text": {"type": "string"},
                                        "best_match_original_run": {"type": "integer"}
                                    },
                                    "required": ["text", "best_match_original_run"]
                                }
                            }
                        },
                        "required": ["id", "translated", "translated_runs"]
                    }
                }
            },
            "required": ["translations"]
        }
        
        self.genai_client = genai.Client()
        self.model_name = model_name

        self.lang_map = {
            "ja": "Japanese",
            "en": "English",
            "ko": "Korean",
            "zh": "Chinese",
            "es": "Spanish",
            "fr": "French",
            "de": "German"
        }
        assert self.source_lang in self.lang_map, f"Unsupported source language: {self.source_lang}"
        assert self.target_lang in self.lang_map, f"Unsupported target language: {self.target_lang}"

    def extract_all_texts_with_positions(self, pptx_path: str) -> PresentationData:
        """PPTXファイルから全テキストと位置情報を抽出"""
        prs = Presentation(pptx_path)
        all_pairs = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # 再帰的にシェイプを処理（グループ化された図形にも対応）
                self._process_shape_recursive(shape, slide_idx, shape_idx, all_pairs)
        
        return PresentationData(all_pairs=all_pairs)
    
    def _process_shape_recursive(self, shape: SlideShapes, slide_idx: int, shape_idx: int, all_pairs: List[TextPositionPair], group_path: Optional[List[int]] = None) -> None:
        """シェイプを再帰的に処理（グループ化された図形にも対応）"""
        if group_path is None:
            group_path = []
        
        # グループ化された図形の場合
        if hasattr(shape, 'shapes'):
            # グループ内の各シェイプを再帰的に処理
            for sub_shape_idx, sub_shape in enumerate(shape.shapes):
                new_group_path = group_path + [sub_shape_idx]
                self._process_shape_recursive(sub_shape, slide_idx, shape_idx, all_pairs, new_group_path)
        
        # 通常のテキストフレーム
        elif hasattr(shape, 'text_frame') and shape.text_frame:
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # runのスタイル情報を収集
                run_styles = self._extract_run_styles(paragraph)
                
                # グループ化されている場合の位置情報
                if group_path:
                    position = GroupedShapePosition(
                        slide_idx=slide_idx,
                        shape_idx=shape_idx,
                        para_idx=para_idx,
                        original_text=text,
                        run_styles=run_styles,
                        group_path=group_path.copy()
                    )
                else:
                    position = TextPosition(
                        slide_idx=slide_idx,
                        shape_idx=shape_idx,
                        para_idx=para_idx,
                        original_text=text,
                        run_styles=run_styles
                    )
                
                pair = TextPositionPair(text=text, position=position)
                all_pairs.append(pair)
        
        # テーブル
        elif hasattr(shape, 'table') and shape.table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if not (hasattr(cell, 'text_frame') and cell.text_frame):
                        continue
                    for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        
                        # runのスタイル情報を収集
                        run_styles = self._extract_run_styles(paragraph)
                        
                        # テーブルセル用の位置情報（グループ化されている場合も考慮）
                        if group_path:
                            # グループ化されたテーブルの場合は、今回は対応しない
                            # 必要に応じて後で実装
                            continue
                        else:
                            position = TableCellPosition(
                                slide_idx=slide_idx,
                                shape_idx=shape_idx,
                                table_row=row_idx,
                                table_col=col_idx,
                                para_idx=para_idx,
                                original_text=text,
                                run_styles=run_styles
                            )
                        
                        pair = TextPositionPair(text=text, position=position)
                        all_pairs.append(pair)
    
    def _extract_run_styles(self, paragraph: _Paragraph) -> List[RunStyleInfo]:
        """paragraphからrunのスタイル情報を抽出"""
        run_styles = []
        for run in paragraph.runs:
            if not run.text:
                continue
            
            font = run.font
            
            # フォントサイズの取得
            font_size_pt = None
            if hasattr(font, 'size') and font.size is not None:
                font_size_pt = font.size.pt
            
            # 色情報の取得
            color_rgb = None
            color_theme = None
            color_brightness = None
            if hasattr(font, 'color') and font.color:
                color_obj = font.color
                if hasattr(color_obj, 'rgb') and color_obj.rgb:
                    color_rgb = color_obj.rgb
                elif hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                    color_theme = int(color_obj.theme_color)
                    if hasattr(color_obj, 'brightness') and color_obj.brightness is not None:
                        color_brightness = float(color_obj.brightness)
            
            run_style = RunStyleInfo(
                text=run.text,
                font_name=getattr(font, 'name', None),
                font_size=font_size_pt,
                bold=getattr(font, 'bold', None),
                italic=getattr(font, 'italic', None),
                underline=getattr(font, 'underline', None),
                color_rgb=color_rgb,
                color_theme=color_theme,
                color_brightness=color_brightness
            )
            run_styles.append(run_style)
        
        return run_styles
    
    def translate_presentation_with_gemini(self, presentation_data: PresentationData) -> Tuple[List[str], List[List[Dict[str, Any]]]]:
        """Gemini APIを使用したプレゼンテーション全体の翻訳"""
        if not presentation_data.all_pairs:
            return [], []
        
        texts = [pair.text for pair in presentation_data.all_pairs]
        self.logger.info(f"翻訳対象テキスト数: {len(texts)}")

        source_lang_name = self.lang_map.get(self.source_lang, self.source_lang)
        target_lang_name = self.lang_map.get(self.target_lang, self.target_lang)
        
        # 翻訳対象テキストとrun情報を整形
        texts_data = []
        for i, pair in enumerate(presentation_data.all_pairs):
            runs_info = []
            for j, run_style in enumerate(pair.position.run_styles):
                runs_info.append({
                    "index": j,
                    "text": run_style.text
                })
            
            text_data = {
                "id": i,
                "text": pair.text,
                "runs": runs_info
            }
            texts_data.append(text_data)
        
        texts_json = json.dumps(texts_data, ensure_ascii=False, indent=2)
        
        prompt = f"""
Translate the following texts from {source_lang_name} to {target_lang_name}.

IMPORTANT REQUIREMENTS:
- Translate each text as a coherent paragraph/sentence
- Maintain the same tone and style throughout
- Keep technical terms consistent across all translations  
- Preserve the original text's nuance, emotional tone, and formality level

CRITICAL - TRANSLATED RUN SEGMENTATION:
After translating each text, you must:

1. Split the translated text into meaningful runs (segments) that correspond to semantic units
2. For each translated run, identify which original run it is most semantically similar to
3. The number of translated runs may differ from original runs - focus on semantic meaning

SEGMENTATION PRINCIPLES:
- Split translated text at natural semantic boundaries (words, phrases, concepts)
- Consider which parts should have similar formatting (emphasis, highlighting, etc.)
- Each translated run should correspond to the most semantically similar original run
- Preserve formatting intent: if original run was emphasized, find the equivalent emphasis in translation

EXAMPLES:
Original: "Hello **world**" (runs: ["Hello ", "world"])
Translation: "こんにちは世界" 
Translated runs: [
  {{"text": "こんにちは", "best_match_original_run": 0}},
  {{"text": "世界", "best_match_original_run": 1}}
]

Original: "**Important:** Please read" (runs: ["Important", ":", " Please read"])
Translation: "重要：必ずお読みください"
Translated runs: [
  {{"text": "重要", "best_match_original_run": 0}},
  {{"text": "：", "best_match_original_run": 1}},
  {{"text": "必ずお読みください", "best_match_original_run": 2}}
]

Input texts with runs:
{texts_json}

For each translation, provide:
1. "translated": The complete translated text
2. "translated_runs": Array of translated runs with their best matching original run indices
"""
        
        self.logger.info("翻訳リクエストを送信中...")
        response = self.genai_client.models.generate_content(
            model=self.model_name,
            contents=prompt,
            config={
                "response_mime_type": "application/json",
                "response_schema": self.translation_schema,
            }
        )

        # JSONレスポンスを解析
        translations, translated_runs = self._parse_json_translations_with_runs(response.text, presentation_data.all_pairs)
        return translations, translated_runs
    
    def _parse_json_translations_with_runs(self, response_text: str, all_pairs: List[TextPositionPair]) -> Tuple[List[str], List[List[Dict[str, Any]]]]:
        """JSON形式の翻訳レスポンス（翻訳後runs付き）をパース"""
        try:
            response_data: Dict[str, Any] = json.loads(response_text)
        except json.JSONDecodeError as e:
            self.logger.error(f"JSONパースエラー: {e}")
            raise ValueError("Invalid JSON response format")
        translations_data = response_data.get("translations", [])
        
        # IDでソートして順序を保証
        translations_data.sort(key=lambda x: x.get("id", 0))

        translations = []
        translated_runs_list = []

        for i, pair in enumerate(all_pairs):
            found_translation = None
            found_translated_runs = []

            for trans_item in translations_data:
                if trans_item.get("id") == i:
                    found_translation = trans_item.get("translated", "")
                    found_translated_runs = trans_item.get("translated_runs", [])
                    break
            
            if found_translation is None:
                # デフォルトの処理：翻訳前のrunsをそのまま使用
                found_translation = pair.text
                found_translated_runs = [
                    {"text": style.text, "best_match_original_run": j}
                    for j, style in enumerate(pair.position.run_styles)
                ]
            
            translations.append(found_translation)
            translated_runs_list.append(found_translated_runs)
        
        self.logger.info(f"翻訳完了: {len(translations)} 件")
        return translations, translated_runs_list

    
    def _create_fallback_mapping(self, translated_text: str, original_run_styles: List[RunStyleInfo]) -> List[Dict[str, Any]]:
        """フォールバックのrun mappingを作成"""
        if not original_run_styles:
            return [{"original_run_index": 0, "translated_text": translated_text}]
        
        # 最初のrunに全テキストを割り当て
        return [{"original_run_index": 0, "translated_text": translated_text}]
    
    def apply_run_style(self, run: _Run, run_style: RunStyleInfo) -> None:
        """runにスタイルを適用"""
        try:
            font = run.font
            
            # フォント名
            if run_style.font_name:
                font.name = run_style.font_name
            
            # フォントサイズ
            if run_style.font_size:
                from pptx.util import Pt
                font.size = Pt(run_style.font_size)
            
            # 太字・斜体・下線
            if run_style.bold is not None:
                font.bold = run_style.bold
            if run_style.italic is not None:
                font.italic = run_style.italic
            if run_style.underline is not None:
                font.underline = run_style.underline
            
            # 色の適用
            if run_style.color_rgb is not None:
                from pptx.dml.color import RGBColor
                r, g, b = run_style.color_rgb
                font.color.rgb = RGBColor(r, g, b)
            elif run_style.color_theme is not None:
                from pptx.enum.dml import MSO_THEME_COLOR
                theme_colors = {
                    0: MSO_THEME_COLOR.BACKGROUND_1,
                    1: MSO_THEME_COLOR.TEXT_1,
                    2: MSO_THEME_COLOR.BACKGROUND_2,
                    3: MSO_THEME_COLOR.TEXT_2,
                    4: MSO_THEME_COLOR.ACCENT_1,
                    5: MSO_THEME_COLOR.ACCENT_2,
                    6: MSO_THEME_COLOR.ACCENT_3,
                    7: MSO_THEME_COLOR.ACCENT_4,
                    8: MSO_THEME_COLOR.ACCENT_5,
                    9: MSO_THEME_COLOR.ACCENT_6,
                    10: MSO_THEME_COLOR.HYPERLINK,
                    11: MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
                    12: MSO_THEME_COLOR.DARK_1,
                    13: MSO_THEME_COLOR.LIGHT_1,
                    14: MSO_THEME_COLOR.DARK_2,
                    15: MSO_THEME_COLOR.LIGHT_2
                }
                if run_style.color_theme in theme_colors:
                    font.color.theme_color = theme_colors[run_style.color_theme]
                    if run_style.color_brightness is not None:
                        font.color.brightness = run_style.color_brightness
                        
        except Exception as e:
            self.logger.debug(f"スタイル適用エラー: {e}")
    
    def replace_text_with_translated_runs(self, paragraph: _Paragraph, translated_text: str, original_run_styles: List[RunStyleInfo], translated_runs: List[Dict[str, Any]]) -> None:
        """翻訳後runsに基づいて最適なスタイルを適用してテキストを置換"""
        try:
            if not original_run_styles or not translated_runs:
                # スタイル情報や翻訳runs情報がない場合は単純置換
                paragraph.text = translated_text
                return
            
            # paragraphをクリア
            paragraph.text = ""
            
            # 翻訳後runsに基づいて新しいrunを作成
            for translated_run in translated_runs:
                run_text = translated_run.get("text", "")
                best_match_original_index = translated_run.get("best_match_original_run", 0)
                
                if not run_text:
                    continue
                
                # 対応する元のスタイルを取得
                if 0 <= best_match_original_index < len(original_run_styles):
                    original_style = original_run_styles[best_match_original_index]
                else:
                    # インデックスが範囲外の場合は最初のスタイルを使用
                    original_style = original_run_styles[0] if original_run_styles else None
                
                # 新しいrunを追加
                run = paragraph.add_run()
                run.text = run_text
                
                # 意味的に最も近い元のrunのスタイルを適用
                if original_style:
                    self.apply_run_style(run, original_style)
        
        except Exception as e:
            # エラーの場合は単純置換にフォールバック
            self.logger.error(f"翻訳runs置換エラー: {e}")
            paragraph.text = translated_text
    
    def translate_pptx(self, input_path: str, output_path: str) -> bool:
        """PPTXファイルを翻訳"""
        try:    
            # ファイルをコピー
            self.logger.info(f"ファイルをコピー中: {input_path} -> {output_path}")
            shutil.copy2(input_path, output_path)
            
            # 全テキストと位置情報を抽出
            presentation_data = self.extract_all_texts_with_positions(input_path)
            
            if not presentation_data.all_pairs:
                self.logger.warning("翻訳対象のテキストが見つかりません")
                return False

            self.logger.info(f"翻訳開始: {input_path}")

            # プレゼンテーション全体を翻訳
            translations, translated_runs_list = self.translate_presentation_with_gemini(presentation_data)

            # 翻訳対応をログ出力
            if self.show_correspondence:
                self._log_correspondence(presentation_data, translations)

            # コピーしたファイルを開いて翻訳を適用
            prs = Presentation(output_path)
            
            for i, (translation, translated_runs, pair) in enumerate(zip(translations, translated_runs_list, presentation_data.all_pairs)):
                try:
                    position = pair.position
                    
                    # 対象のシェイプと段落を取得
                    slide = prs.slides[position.slide_idx]
                    shape = slide.shapes[position.shape_idx]
                    
                    # グループ化された図形の場合
                    if isinstance(position, GroupedShapePosition):
                        # group_pathに従って階層的にシェイプを取得
                        current_shape = shape
                        for group_idx in position.group_path:
                            current_shape = current_shape.shapes[group_idx]
                        paragraph = current_shape.text_frame.paragraphs[position.para_idx]
                    
                    # テーブルセルの場合
                    elif isinstance(position, TableCellPosition):
                        table = shape.table
                        cell = table.rows[position.table_row].cells[position.table_col]
                        paragraph = cell.text_frame.paragraphs[position.para_idx]
                    
                    else:
                        # 通常のテキストフレーム
                        paragraph = shape.text_frame.paragraphs[position.para_idx]
                    
                    # 翻訳テキストを意味的対応関係に基づいてスタイル適用
                    self.replace_text_with_translated_runs(paragraph, translation, position.run_styles, translated_runs)
                    
                except Exception as e:
                    self.logger.error(f"テキスト置換エラー (インデックス {i}): {e}")
                    continue

            prs.save(output_path)
            self.logger.info(f"{output_path}に保存しました。")
            return True
            
        except Exception as e:
            self.logger.error(f"翻訳処理エラー: {e}")
            return False

    def _log_correspondence(self, presentation_data: PresentationData, translations: List[str]) -> None:
        """翻訳対応をログで表示"""
        self.logger.info("=" * 80)
        self.logger.info("翻訳対応表")
        self.logger.info("=" * 80)
        
        for i, (pair, translation) in enumerate(zip(presentation_data.all_pairs, translations)):
            original_text = pair.text
            position = pair.position
            
            # 位置情報の文字列表現を作成
            if isinstance(position, TableCellPosition):
                location = f"スライド{position.slide_idx + 1}, 表 行{position.table_row + 1}列{position.table_col + 1}"
            elif isinstance(position, GroupedShapePosition):
                location = f"スライド{position.slide_idx + 1}, グループ図形 {'.'.join(map(str, position.group_path))}"
            else:
                location = f"スライド{position.slide_idx + 1}, 図形{position.shape_idx + 1}"
            
            self.logger.info(f"[{i + 1}] {location}")
            self.logger.info(f"元の文: {original_text}")
            self.logger.info(f"翻訳文: {translation}")
            self.logger.info("-" * 40)
        
        self.logger.info("=" * 80)


def main() -> None:
    parser = argparse.ArgumentParser(description="PPTX翻訳スクリプト (Gemini API使用)")
    parser.add_argument("input_file", help="入力PPTXファイル")
    parser.add_argument("-o", "--output", help="出力PPTXファイル (デフォルト: <input_file_name>_<target_lang>.pptx)")
    parser.add_argument("-s", "--source", default="ja", help="翻訳元言語 (デフォルト: ja)")
    parser.add_argument("-t", "--target", default="en", help="翻訳先言語 (デフォルト: en)")
    parser.add_argument("-m", "--model", default="gemini-2.5-flash", help="使用するモデル名 (デフォルト: gemini-2.5-flash)")
    parser.add_argument("--show-correspondence", action="store_true", help="翻訳前後の対応をログで表示")

    args = parser.parse_args()
    
    # 入力ファイルの確認
    if not Path(args.input_file).exists():
        print(f"エラー: ファイルが見つかりません: {args.input_file}")
        sys.exit(1)
    
    # 出力ファイル名の生成
    if args.output:
        output_file = args.output
    else:
        input_path = Path(args.input_file)
        output_file = str(input_path.parent / f"{input_path.stem}_{args.target}.pptx")

    # 翻訳処理
    translator = PPTXTranslator(
        args.source, 
        args.target, 
        args.model,
        show_correspondence=args.show_correspondence
    )
    success = translator.translate_pptx(args.input_file, output_file)
    
    if success:
        print("翻訳が完了しました。")
    else:
        print("翻訳に失敗しました。")
        sys.exit(1)


if __name__ == "__main__":
    main()